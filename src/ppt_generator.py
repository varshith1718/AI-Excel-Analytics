from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from datetime import datetime
import os


class PPTGenerator:

    def __init__(self):

        self.output_folder = "output"
        os.makedirs(self.output_folder, exist_ok=True)

        self.prs = Presentation()

        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

        self.white = RGBColor(255, 255, 255)
        self.black = RGBColor(35, 35, 35)

        self.dark_blue = RGBColor(30, 58, 138)
        self.blue = RGBColor(37, 99, 235)
        self.sky = RGBColor(96, 165, 250)

        self.green = RGBColor(34, 197, 94)
        self.orange = RGBColor(245, 158, 11)
        self.red = RGBColor(239, 68, 68)

        self.light = RGBColor(245, 248, 252)

    def set_background(self, slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.white

    def add_title(self, slide, title):
        box = slide.shapes.add_textbox(
            Inches(0.4),
            Inches(0.2),
            Inches(8),
            Inches(0.6)
        )
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(26)
        p.font.color.rgb = self.dark_blue

    def add_cover_slide(self):
        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]
        )

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.dark_blue

        title = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(1.2),
            Inches(11),
            Inches(1)
        )

        p = title.text_frame.paragraphs[0]
        p.text = "AI EXCEL ANALYTICS REPORT"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

        p = title.text_frame.add_paragraph()
        p.text = "Executive Dashboard"
        p.font.size = Pt(20)
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

        p = title.text_frame.add_paragraph()
        p.text = datetime.now().strftime("%d %B %Y")
        p.font.size = Pt(14)
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

        footer = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(6.7),
            Inches(11),
            Inches(0.3)
        )

        p = footer.text_frame.paragraphs[0]
        p.text = "Generated using AI Excel Analytics & Report Generator"
        p.font.size = Pt(12)
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

    def add_card(self, slide, left, top, title, value, color):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            top,
            Inches(2.1),
            Inches(1)
        )

        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = color

        tf = card.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = str(value)
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

    def add_dashboard_slide(self, dataset_info, score, recommendations):
        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]
        )

        self.set_background(slide)
        self.add_title(slide, "Executive Dashboard")

        self.add_card(
            slide, Inches(0.4), Inches(0.9),
            "ROWS", dataset_info["Rows"], self.blue
        )

        self.add_card(
            slide, Inches(2.7), Inches(0.9),
            "COLUMNS", dataset_info["Columns"], self.sky
        )

        self.add_card(
            slide, Inches(5.0), Inches(0.9),
            "MISSING", dataset_info["Missing Values"], self.orange
        )

        self.add_card(
            slide, Inches(7.3), Inches(0.9),
            "DUPLICATES", dataset_info["Duplicate Rows"], self.red
        )

        self.add_card(
            slide, Inches(9.6), Inches(0.9),
            "SCORE", f"{score}/100", self.green
        )

        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(2.2),
            Inches(12.2),
            Inches(3.9)
        )

        box.fill.solid()
        box.fill.fore_color.rgb = self.light
        box.line.color.rgb = self.sky

        tf = box.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = "Executive Summary"
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = self.dark_blue

        summary = [
            f"Dataset contains {dataset_info['Rows']} records.",
            f"Dataset has {dataset_info['Columns']} columns.",
            f"Missing Values : {dataset_info['Missing Values']}",
            f"Duplicate Rows : {dataset_info['Duplicate Rows']}",
            f"Overall Dataset Score : {score}/100"
        ]

        for item in summary:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(16)
            p.font.color.rgb = self.black

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "Recommendations"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = self.dark_blue

        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = "✓ " + rec
            p.font.size = Pt(15)
            p.font.color.rgb = self.black

    def add_dataset_slide(self, dataset_info):
        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]
        )

        self.set_background(slide)
        self.add_title(slide, "Dataset Overview")

        table = slide.shapes.add_table(
            8, 2,
            Inches(0.6), Inches(0.9),
            Inches(6), Inches(3.8)
        ).table

        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"

        headers = [
            ("Rows", dataset_info["Rows"]),
            ("Columns", dataset_info["Columns"]),
            ("Missing Values", dataset_info["Missing Values"]),
            ("Duplicate Rows", dataset_info["Duplicate Rows"]),
            ("Numeric Columns", len(dataset_info["Numeric Columns"])),
            ("Categorical Columns", len(dataset_info["Categorical Columns"])),
            ("Column Names", len(dataset_info["Column Names"]))
        ]

        for i, (key, value) in enumerate(headers, start=1):
            table.cell(i, 0).text = str(key)
            table.cell(i, 1).text = str(value)

        for row in table.rows:
            for cell in row.cells:
                cell.text_frame.paragraphs[0].font.size = Pt(14)

        for col in range(2):
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.blue
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.color.rgb = self.white

        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(7),
            Inches(0.9),
            Inches(5.5),
            Inches(5)
        )

        box.fill.solid()
        box.fill.fore_color.rgb = self.light
        box.line.color.rgb = self.sky

        tf = box.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = "Columns"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = self.dark_blue

        for col in dataset_info["Column Names"]:
            p = tf.add_paragraph()
            p.text = "• " + str(col)
            p.font.size = Pt(13)
            p.font.color.rgb = self.black

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "Data Quality"
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = self.dark_blue

        quality = "Excellent"

        if dataset_info["Missing Values"] > 0:
            quality = "Good"

        if dataset_info["Duplicate Rows"] > 0:
            quality = "Needs Cleaning"

        p = tf.add_paragraph()
        p.text = f"Overall Quality : {quality}"
        p.font.size = Pt(15)
        p.font.bold = True

        if quality == "Excellent":
            p.font.color.rgb = self.green
        elif quality == "Good":
            p.font.color.rgb = self.orange
        else:
            p.font.color.rgb = self.red

    def add_chart_slide(self, charts):
        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]
        )

        self.set_background(slide)
        self.add_title(slide, "Data Visualization")

        chart_names = [
            ("Bar", 0.5, 0.9),
            ("Pie", 6.8, 0.9),
            ("Histogram", 0.5, 3.8),
            ("Line", 6.8, 3.8)
        ]

        for name, left, top in chart_names:

            if name not in charts:
                continue

            image = charts[name]

            if image is None:
                continue

            if not os.path.exists(image):
                continue

            title = slide.shapes.add_textbox(
                Inches(left),
                Inches(top),
                Inches(2),
                Inches(0.3)
            )

            p = title.text_frame.paragraphs[0]
            p.text = name + " Chart"
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = self.dark_blue

            slide.shapes.add_picture(
                image,
                Inches(left),
                Inches(top + 0.3),
                width=Inches(5.5),
                height=Inches(2.5)
            )

    def add_ai_insights_slide(self, insights, risks, recommendations, score):
        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]
        )

        self.set_background(slide)
        self.add_title(slide, "AI Executive Insights")

        score_color = self.green

        if score < 80:
            score_color = self.orange

        if score < 60:
            score_color = self.red

        self.add_card(
            slide, Inches(9.8), Inches(0.8),
            "DATA SCORE", f"{score}/100", score_color
        )

        insight_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(0.9),
            Inches(8.8),
            Inches(2.3)
        )

        insight_box.fill.solid()
        insight_box.fill.fore_color.rgb = self.light
        insight_box.line.color.rgb = self.sky

        tf = insight_box.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = "AI Insights"
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = self.dark_blue

        for item in insights:
            p = tf.add_paragraph()
            p.text = "• " + str(item)
            p.font.size = Pt(15)
            p.font.color.rgb = self.black

        risk_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(3.5),
            Inches(6),
            Inches(2.8)
        )

        risk_box.fill.solid()
        risk_box.fill.fore_color.rgb = self.light
        risk_box.line.color.rgb = self.orange

        tf = risk_box.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = "Risk Analysis"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = self.dark_blue

        for risk in risks:
            p = tf.add_paragraph()
            p.text = "⚠ " + str(risk)
            p.font.size = Pt(14)
            p.font.color.rgb = self.black

        rec_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(6.8),
            Inches(3.5),
            Inches(5.5),
            Inches(2.8)
        )

        rec_box.fill.solid()
        rec_box.fill.fore_color.rgb = self.light
        rec_box.line.color.rgb = self.green

        tf = rec_box.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = "Recommendations"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = self.dark_blue

        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = "✓ " + str(rec)
            p.font.size = Pt(14)
            p.font.color.rgb = self.black

        footer = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(6.55),
            Inches(11.8),
            Inches(0.35)
        )

        footer.fill.solid()
        footer.fill.fore_color.rgb = self.blue
        footer.line.color.rgb = self.blue

        tf = footer.text_frame

        p = tf.paragraphs[0]
        p.text = (
            "AI-generated insights based on dataset quality, "
            "missing values, duplicate records and statistical analysis."
        )
        p.font.size = Pt(12)
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

    def add_recommendation_slide(self, recommendations):
        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]
        )

        self.set_background(slide)
        self.add_title(slide, "Executive Recommendations")

        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(0.9),
            Inches(12),
            Inches(5.8)
        )

        box.fill.solid()
        box.fill.fore_color.rgb = self.light
        box.line.color.rgb = self.sky

        tf = box.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = "Recommended Actions"
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = self.dark_blue

        if len(recommendations) == 0:
            recommendations = [
                "Dataset quality is good.",
                "Continue monitoring data quality.",
                "Generate reports periodically."
            ]

        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = "✓ " + str(rec)
            p.font.size = Pt(16)
            p.font.color.rgb = self.black

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "Best Practices"
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = self.dark_blue

        best_practices = [
            "Validate data before analysis.",
            "Handle missing values carefully.",
            "Remove duplicate records.",
            "Maintain consistent column names.",
            "Review data quality regularly.",
            "Generate dashboards for business users."
        ]

        for item in best_practices:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(15)
            p.font.color.rgb = self.black

        footer = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(6.55),
            Inches(12),
            Inches(0.35)
        )

        footer.fill.solid()
        footer.fill.fore_color.rgb = self.green
        footer.line.color.rgb = self.green

        tf = footer.text_frame

        p = tf.paragraphs[0]
        p.text = (
            "Following these recommendations will improve data quality and analytics accuracy."
        )
        p.font.size = Pt(12)
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

    def add_executive_summary_slide(self, dataset_info, score):

        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # ---------------- Title ----------------
        title = slide.shapes.add_textbox(
            Inches(0.6),
            Inches(0.3),
            Inches(8),
            Inches(0.5)
        )

        p = title.text_frame.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 78, 121)

        # ---------------- Rounded Box ----------------
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.7),
            Inches(1.0),
            Inches(8.2),
            Inches(5.2)
        )

        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 248, 255)

        box.line.color.rgb = RGBColor(70, 130, 255)

        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True

        rows = dataset_info.get("Rows", 0)
        cols = dataset_info.get("Columns", 0)
        missing = dataset_info.get("Missing Values", 0)
        duplicates = dataset_info.get("Duplicate Rows", 0)

        quality = "Excellent"

        if missing > 0 or duplicates > 0:
            quality = "Good"

        if missing > 20:
            quality = "Needs Improvement"

        p = tf.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(22)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(31, 78, 121)

        p = tf.add_paragraph()
        p.text = ""
        p.font.size = Pt(8)

        p = tf.add_paragraph()
        p.text = "The uploaded Excel file has been analyzed successfully."
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)

        p = tf.add_paragraph()
        p.text = f"The dataset contains {rows} records and {cols} columns."
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)

        p = tf.add_paragraph()
        p.text = f"Overall data quality is {quality}."
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)

        p = tf.add_paragraph()
        p.text = ""
        p.font.size = Pt(8)

        p = tf.add_paragraph()
        p.text = "Dataset Score"
        p.font.bold = True
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(31, 78, 121)

        p = tf.add_paragraph()
        p.text = f"{score}/100"
        p.font.bold = True
        p.font.size = Pt(30)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(0, 176, 80)

    def add_thankyou_slide(self):
        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]
        )

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.dark_blue

        title = slide.shapes.add_textbox(
            Inches(1),
            Inches(1.6),
            Inches(11),
            Inches(1)
        )

        p = title.text_frame.paragraphs[0]
        p.text = "THANK YOU"
        p.font.bold = True
        p.font.size = Pt(34)
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

        p = title.text_frame.add_paragraph()
        p.text = "AI Excel Analytics & Report Generator"
        p.font.size = Pt(22)
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

        p = title.text_frame.add_paragraph()
        p.text = "Generated Automatically using Python"
        p.font.size = Pt(16)
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

        footer = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(6.7),
            Inches(12),
            Inches(0.3)
        )

        p = footer.text_frame.paragraphs[0]
        p.text = datetime.now().strftime("Generated on %d %B %Y")
        p.font.size = Pt(12)
        p.font.color.rgb = self.white
        p.alignment = PP_ALIGN.CENTER

    def generate_presentation(
        self,
        dataset_info,
        charts,
        insights,
        risks,
        recommendations,
        score
    ):
        self.add_cover_slide()

        self.add_dashboard_slide(
            dataset_info,
            score,
            recommendations
        )

        self.add_dataset_slide(dataset_info)

        self.add_chart_slide(charts)

        self.add_ai_insights_slide(
            insights,
            risks,
            recommendations,
            score
        )

        self.add_recommendation_slide(recommendations)

        self.add_executive_summary_slide(
            dataset_info,
            score
        )

        self.add_thankyou_slide()
        filename = datetime.now().strftime("%Y%m%d_%H%M%S")

        output_file = os.path.join(
            self.output_folder,
            f"Excel_Analytics_Report_{filename}.pptx"
        )

        self.prs.save(output_file)

        return output_file